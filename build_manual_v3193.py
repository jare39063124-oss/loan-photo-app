#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
资产盘点拍照工具 使用说明书 v3.22.24
从零重建 PPTX —— 简约明亮配色，修复排版问题
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ============================================================
# 设计系统
# ============================================================
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 配色（简约明亮）
C_BG       = RGBColor(0xFF, 0xFF, 0xFF)   # 白底
C_BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)   # 浅灰蓝底
C_ACCENT   = RGBColor(0x21, 0x96, 0xF3)   # 活力蓝
C_ACCENT_D = RGBColor(0x14, 0x65, 0xC7)   # 深蓝
C_SUCCESS  = RGBColor(0x33, 0xB3, 0x5C)   # 绿
C_DANGER   = RGBColor(0xF4, 0x43, 0x36)   # 红
C_WARNING  = RGBColor(0xFF, 0x9E, 0x0A)  # 琥珀
C_TEXT     = RGBColor(0x1F, 0x21, 0x26)   # 近黑
C_TEXT_DIM = RGBColor(0x6B, 0x73, 0x80)   # 中灰
C_BORDER   = RGBColor(0xDC, 0xE0, 0xE8)   # 浅灰边框
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD_BG  = RGBColor(0xFA, 0xFB, 0xFC)

FONT = '微软雅黑'
IMG_DIR = r'd:\hermes\loan_photo_app\ppt_images'
DST = r'C:\Users\Administrator\Desktop\资产盘点拍照工具-使用说明书-v3.22.24.pptx'

# ============================================================
# 辅助函数
# ============================================================
def add_slide(prs):
    """添加空白幻灯片"""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    # 白色背景
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_BG
    return slide

def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=None, radius=None):
    """添加矩形/圆角矩形"""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    if radius is not None:
        # 调整圆角半径
        shape.adjustments[0] = radius
    return shape

def add_text(slide, l, t, w, h, text, font_size=14, bold=False, color=C_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name=FONT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        if i == 0:
            para = p
        else:
            para = tf.add_paragraph()
            para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txBox

def add_text_in_shape(shape, text, font_size=14, bold=False, color=C_WHITE,
                      align=PP_ALIGN.CENTER, font_name=FONT, word_wrap=True):
    """在形状内添加文字"""
    tf = shape.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    tf.paragraphs[0].alignment = align
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name

def add_title_bar(slide, chapter_num, title):
    """添加幻灯片顶部标题栏"""
    # 顶部蓝色条
    bar = add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.75), C_ACCENT)
    # 章节编号
    add_text(slide, Inches(0.5), Inches(0.05), Inches(10), Inches(0.65),
             "第%d章  %s" % (chapter_num, title), font_size=24, bold=True, color=C_WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    # 底部细线
    add_rect(slide, Inches(0), Inches(0.75), SLIDE_W, Pt(2), C_ACCENT_D)

def add_page_num(slide, num, total=16):
    """添加页码"""
    add_text(slide, Inches(12.3), Inches(7.0), Inches(0.9), Inches(0.35),
             "%d / %d" % (num, total), font_size=10, color=C_TEXT_DIM,
             align=PP_ALIGN.RIGHT)

def add_card(slide, l, t, w, h, title=None, fill=C_CARD_BG):
    """添加卡片容器"""
    card = add_rect(slide, l, t, w, h, fill, C_BORDER, Pt(1), radius=0.03)
    if title:
        add_text(slide, l + Inches(0.2), t + Inches(0.12), w - Inches(0.4), Inches(0.45),
                 title, font_size=18, bold=True, color=C_ACCENT_D)
    return card

def add_number_badge(slide, l, t, num, size=Inches(0.55)):
    """添加章节编号徽章（足够宽，支持两位数不竖排）"""
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    badge.fill.solid()
    badge.fill.fore_color.rgb = C_ACCENT
    badge.line.fill.background()
    add_text_in_shape(badge, str(num), font_size=14, bold=True, color=C_WHITE, word_wrap=False)
    return badge

def add_bullet_list(slide, l, t, w, h, items, font_size=15, spacing=Pt(8)):
    """添加带●符号的列表"""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = spacing
        run = p.add_run()
        run.text = "●  " + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = C_TEXT
        run.font.name = FONT
    return txBox

# ============================================================
# GUI 模拟图辅助函数
# ============================================================
def add_phone_mockup(slide, l, t, w, h, title="资产盘点拍照"):
    """添加手机界面模拟外框"""
    # 外框
    frame = add_rect(slide, l, t, w, h, C_WHITE, C_BORDER, Pt(1.5), radius=0.04)
    # 顶部标题栏
    bar_h = Inches(0.5)
    bar = add_rect(slide, l, t, w, bar_h, C_ACCENT, radius=0.04)
    add_text(slide, l + Inches(0.1), t + Inches(0.05), w - Inches(0.2), bar_h - Inches(0.1),
             title, font_size=11, bold=True, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)
    return frame

# ============================================================
# 创建演示文稿
# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ============================================================
# Slide 1: 封面
# ============================================================
slide = add_slide(prs)
# 左侧蓝色面板
left_panel = add_rect(slide, Inches(0), Inches(0), Inches(6.5), SLIDE_H, C_ACCENT)
# 标题
add_text(slide, Inches(0.8), Inches(1.8), Inches(5.2), Inches(1.2),
         "资产盘点专项拍照工具", font_size=36, bold=True, color=C_WHITE)
add_text(slide, Inches(0.8), Inches(3.0), Inches(5.2), Inches(0.6),
         "使用说明书", font_size=24, color=C_WHITE)
# 版本
add_text(slide, Inches(0.8), Inches(3.8), Inches(5.2), Inches(0.5),
         "v3.22.24", font_size=18, color=C_WHITE)
# 底部署名
add_text(slide, Inches(0.8), Inches(5.5), Inches(5.2), Inches(0.8),
         ["问题咨询电话：15940454123", "2026年6月"], font_size=14, color=C_WHITE)
# 右侧装饰区（替换损坏的生成图，使用矢量图形设计）
right_panel = add_rect(slide, Inches(6.5), Inches(0), Inches(6.83), SLIDE_H, C_BG_LIGHT)
# 大菱形装饰
diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(8.5), Inches(1.5), Inches(3.0), Inches(3.0))
diamond.fill.solid()
diamond.fill.fore_color.rgb = C_ACCENT_D
diamond.line.fill.background()
# 内圈
inner = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(9.2), Inches(2.2), Inches(1.6), Inches(1.6))
inner.fill.solid()
inner.fill.fore_color.rgb = C_WHITE
inner.line.color.rgb = C_ACCENT
inner.line.width = Pt(3)
# 中心文字
add_text(slide, Inches(9.2), Inches(2.6), Inches(1.6), Inches(0.8),
         "资产\n盘点", font_size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
# 底部装饰条
add_rect(slide, Inches(7.5), Inches(5.5), Inches(5.0), Pt(3), C_ACCENT)
add_text(slide, Inches(7.5), Inches(5.7), Inches(5.0), Inches(0.5),
         "抵押物 / 抵债资产 现场勘查专用", font_size=14, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 2: 目录
# ============================================================
slide = add_slide(prs)
# 标题
add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "目  录", font_size=28, bold=True, color=C_ACCENT_D, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(1.0), Inches(2.3), Pt(3), C_ACCENT)

chapters = [
    (1, "应用概述"), (2, "首次使用"), (3, "Excel格式规范"),
    (4, "客户列表"), (5, "拍照功能"), (6, "照片查看与管理"),
    (7, "设置页面"), (8, "水印说明"), (9, "文件命名规则"),
    (10, "日报表生成"), (11, "AI助手功能"), (12, "日志管理"), (13, "常见问题"),
]

# 两列布局，13 项分 7+6
col_x = [Inches(1.2), Inches(7.0)]
row_y_start = Inches(1.4)
row_h = Inches(0.85)
badge_size = Inches(0.55)

for i, (num, name) in enumerate(chapters):
    col = i // 7  # v3.22.2: 第一列 7 个，第二列 6 个
    row = i % 7
    x = col_x[col]
    y = row_y_start + row * row_h
    # 编号徽章
    add_number_badge(slide, x, y, num, badge_size)
    # 章节名
    add_text(slide, x + badge_size + Inches(0.2), y, Inches(4.5), badge_size,
             "第%d章  %s" % (num, name), font_size=15, color=C_TEXT,
             anchor=MSO_ANCHOR.MIDDLE)

add_page_num(slide, 2)

# ============================================================
# Slide 3: 第一章 应用概述
# ============================================================
slide = add_slide(prs)
add_title_bar(slide, 1, "应用概述")

# 左侧：功能介绍
add_card(slide, Inches(0.4), Inches(1.0), Inches(6.0), Inches(2.6), "功能介绍")
add_bullet_list(slide, Inches(0.6), Inches(1.5), Inches(5.6), Inches(2.0), [
    "Excel批量导入客户地址信息",
    "按客户分类管理资产照片",
    "远景/近景/内部/瑕疵/其他五种拍照类型",
    "自动添加水印（日期、地址、GPS定位）",
    "AI一键生成勘查日报表",
    "AI智能助手（查询拍摄进度）",
], font_size=15)

# 适用场景（不再与功能介绍重叠，独立卡片）
add_card(slide, Inches(0.4), Inches(3.8), Inches(6.0), Inches(1.6), "适用场景")
add_text(slide, Inches(0.6), Inches(4.3), Inches(5.6), Inches(1.0),
         "专项用于抚顺银行资产盘点工作，适用于抵押物、抵债资产现场勘查拍照取证，"
         "支持贷前调查、贷中检查、贷后管理、资产保全等业务场景。",
         font_size=15, color=C_TEXT_DIM)

# 右侧：手机界面模拟
phone_x = Inches(7.2)
phone_y = Inches(1.0)
phone_w = Inches(5.6)
phone_h = Inches(5.8)
add_phone_mockup(slide, phone_x, phone_y, phone_w, phone_h, "资产盘点拍照")
# 模拟标题栏按钮
btn_y = phone_y + Inches(0.6)
add_rect(slide, phone_x + Inches(0.2), btn_y, Inches(1.2), Inches(0.35), C_SUCCESS)
add_text(slide, phone_x + Inches(0.2), btn_y, Inches(1.2), Inches(0.35),
         "AI助手", font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(slide, phone_x + Inches(1.5), btn_y, Inches(1.0), Inches(0.35), C_ACCENT)
add_text(slide, phone_x + Inches(1.5), btn_y, Inches(1.0), Inches(0.35),
         "设置", font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 进度
add_text(slide, phone_x + Inches(3.0), btn_y, Inches(2.4), Inches(0.35),
         "0/0", font_size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
# 工具栏
tool_y = phone_y + Inches(1.05)
add_rect(slide, phone_x + Inches(0.2), tool_y, Inches(2.0), Inches(0.35), C_ACCENT)
add_text(slide, phone_x + Inches(0.2), tool_y, Inches(2.0), Inches(0.35),
         "打开Excel", font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(slide, phone_x + Inches(2.3), tool_y, Inches(1.8), Inches(0.35), C_BORDER)
add_text(slide, phone_x + Inches(2.3), tool_y, Inches(1.8), Inches(0.35),
         "搜索客户名…", font_size=8, color=C_TEXT_DIM, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(slide, phone_x + Inches(4.2), tool_y, Inches(1.2), Inches(0.35), C_ACCENT_D)
add_text(slide, phone_x + Inches(4.2), tool_y, Inches(1.2), Inches(0.35),
         "搜索", font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 客户列表模拟
list_y = phone_y + Inches(1.55)
for i in range(4):
    ry = list_y + Inches(i * 0.7)
    add_rect(slide, phone_x + Inches(0.2), ry, phone_w - Inches(0.4), Inches(0.6),
             C_CARD_BG if i % 2 == 0 else C_WHITE, C_BORDER, Pt(0.5))
    add_text(slide, phone_x + Inches(0.3), ry + Inches(0.05), Inches(1.5), Inches(0.5),
             ["成都投资集团", "沈阳置业有限", "抚顺实业公司", "大连贸易有限"][i],
             font_size=9, bold=True, color=C_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, phone_x + Inches(1.8), ry + Inches(0.05), Inches(1.8), Inches(0.5),
             ["和平区XX街123号", "沈河区YY路45号", "新抚区ZZ道67号", "中山区AA路89号"][i],
             font_size=7, color=C_TEXT_DIM, anchor=MSO_ANCHOR.MIDDLE)
    status_color = C_SUCCESS if i < 2 else C_WARNING
    add_rect(slide, phone_x + Inches(3.7), ry + Inches(0.1), Inches(0.7), Inches(0.35), status_color)
    add_text(slide, phone_x + Inches(3.7), ry + Inches(0.1), Inches(0.7), Inches(0.35),
             ["远2近1", "远1近2", "远0近0", "远0近0"][i],
             font_size=7, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, phone_x + Inches(4.5), ry + Inches(0.1), Inches(0.7), Inches(0.35), C_ACCENT_D)
    add_text(slide, phone_x + Inches(4.5), ry + Inches(0.1), Inches(0.7), Inches(0.35),
             "拍照", font_size=8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 底部日报表按钮
add_rect(slide, phone_x + Inches(0.2), phone_y + Inches(5.0), phone_w - Inches(0.4), Inches(0.5),
         C_SUCCESS, radius=0.1)
add_text(slide, phone_x + Inches(0.2), phone_y + Inches(5.0), phone_w - Inches(0.4), Inches(0.5),
         "AI 一键生成日报表", font_size=12, bold=True, color=C_WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_page_num(slide, 3)

# ============================================================
# Slide 4: 第二章 首次使用
# ============================================================
slide = add_slide(prs)
add_title_bar(slide, 2, "首次使用")

# 左侧：说明文字
add_card(slide, Inches(0.4), Inches(1.0), Inches(6.0), Inches(2.0), "欢迎页")
add_text(slide, Inches(0.6), Inches(1.5), Inches(5.6), Inches(1.4),
         ["首次启动应用时显示欢迎引导页面，简要介绍主要功能。",
          "点击「开始使用」按钮进入主界面。"], font_size=15, color=C_TEXT)

add_card(slide, Inches(0.4), Inches(3.2), Inches(6.0), Inches(3.5), "权限说明")
add_bullet_list(slide, Inches(0.6), Inches(3.7), Inches(5.6), Inches(2.8), [
    "相机权限：用于拍摄资产照片",
    "存储权限：用于保存照片到手机",
    "位置权限：用于获取GPS定位信息（水印中显示经纬度）",
    "网络权限：用于AI助手和日报表生成功能",
    "首次使用时请全部允许",
    "如误拒可在手机设置→应用管理→本应用→权限中重新开启",
], font_size=15)

# 右侧：欢迎页模拟
wx = Inches(7.2)
wy = Inches(1.0)
ww = Inches(5.6)
wh = Inches(5.8)
add_phone_mockup(slide, wx, wy, ww, wh, "欢迎使用")
# Logo
add_text(slide, wx + Inches(1.8), wy + Inches(0.8), Inches(2.0), Inches(0.8),
         "◆", font_size=40, color=C_ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, wx + Inches(0.3), wy + Inches(1.7), ww - Inches(0.6), Inches(0.5),
         "资产盘点专项拍照工具", font_size=15, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
add_text(slide, wx + Inches(0.3), wy + Inches(2.2), ww - Inches(0.6), Inches(0.35),
         "v3.22.24", font_size=10, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)
# 功能列表
features = [
    "四类拍照引导（远景/近景/内部/瑕疵）",
    "水印自选模式（段+位置+字号）",
    "文件命名自选模式（4段下拉）",
    "一键生成勘查日报表",
]
for i, feat in enumerate(features):
    fy = wy + Inches(2.7) + Inches(i * 0.35)
    add_text(slide, wx + Inches(0.5), fy, ww - Inches(1.0), Inches(0.3),
             "●  " + feat, font_size=10, color=C_TEXT)
# 开始按钮
add_rect(slide, wx + Inches(0.8), wy + Inches(4.5), ww - Inches(1.6), Inches(0.55), C_ACCENT, radius=0.15)
add_text(slide, wx + Inches(0.8), wy + Inches(4.5), ww - Inches(1.6), Inches(0.55),
         "开 始 使 用", font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 署名
add_text(slide, wx + Inches(0.3), wy + Inches(5.2), ww - Inches(0.6), Inches(0.35),
         "问题咨询电话：15940454123", font_size=9, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)

add_page_num(slide, 4)

# ============================================================
# Slide 5: 安装说明（v3.22.0 新增）
# ============================================================
slide = add_slide(prs)
add_title_bar(slide, 0, "安装说明")

# 4 步安装流程，2x2 网格布局
step_w = Inches(5.8)
step_h = Inches(2.5)
gap = Inches(0.3)
start_x = Inches(0.4)
start_y = Inches(1.1)

steps = [
    ("1", "微信接收文件",
     "通过微信接收同事发送的 APK 安装文件\n（loan-photo-tool-v3.22.24.apk），\n在微信聊天中点击该文件。"),
    ("2", "打开文件菜单",
     "在微信中打开 APK 文件预览，\n点击右上角「三个点」图标，\n弹出操作菜单。"),
    ("3", "保存到手机",
     "在弹出的对话框中点击\n「保存到手机」或「保存」，\n文件将保存到手机存储中。"),
    ("4", "文件管理器安装",
     "打开手机「文件管理器」，在「最近」\n或「微信」文件夹中找到 APK 文件，\n点击后按提示安装，或使用安装工具打开。"),
]

for i, (num, title, desc) in enumerate(steps):
    col = i % 2
    row = i // 2
    x = start_x + col * (step_w + gap)
    y = start_y + row * (step_h + gap)

    # 卡片背景
    add_card(slide, x, y, step_w, step_h)
    # 编号徽章
    add_number_badge(slide, x + Inches(0.2), y + Inches(0.2), num)
    # 标题
    add_text(slide, x + Inches(0.9), y + Inches(0.15), step_w - Inches(1.1), Inches(0.55),
             title, font_size=18, bold=True, color=C_ACCENT_D, anchor=MSO_ANCHOR.MIDDLE)
    # 说明
    add_text(slide, x + Inches(0.3), y + Inches(0.9), step_w - Inches(0.6), Inches(1.5),
             desc, font_size=13, color=C_TEXT)

# 底部提示
add_rect(slide, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.5), C_BG_LIGHT, C_BORDER, Pt(1), radius=0.05)
add_text(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
         "提示：如安装时提示「未知来源」，请在手机设置→安全→允许安装未知来源应用中开启权限。",
         font_size=13, color=C_TEXT_DIM, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

add_page_num(slide, 5)

# ============================================================
# Slide 6: 第三章 Excel格式规范
# ============================================================
slide = add_slide(prs)
add_title_bar(slide, 3, "Excel格式规范")

add_card(slide, Inches(0.4), Inches(1.0), Inches(12.5), Inches(2.2), "格式要求")
add_bullet_list(slide, Inches(0.6), Inches(1.5), Inches(12.0), Inches(1.6), [
    "仅支持 .xlsx 格式（Excel 2007及以上版本）",
    "第一行为表头行（A1:F1），数据从第二行开始",
    "列顺序固定：A=序号, B=客户名, C=地址概(街道+楼栋), D=地址精确(单元房号), E=性质, F=备注",
    "F列为备注列，由用户手动填写勘查备注，App读取此列内容用于AI生成日报表（严禁编造）",
], font_size=15)

# Excel表格模拟（v3.22.0: 6列格式）
tbl_x = Inches(0.4)
tbl_y = Inches(3.5)
col_widths = [Inches(0.8), Inches(2.2), Inches(3.0), Inches(2.2), Inches(1.5), Inches(2.3)]
headers = ["序号", "客户名", "地址概(街道+楼栋)", "地址精确(单元房号)", "性质", "备注"]
# 表头
cx = tbl_x
for i, (hdr, cw) in enumerate(zip(headers, col_widths)):
    cell = add_rect(slide, cx, tbl_y, cw, Inches(0.45), C_ACCENT)
    add_text(slide, cx, tbl_y, cw, Inches(0.45), hdr,
             font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += cw
# 数据行
data = [
    ["1", "成都投资集团", "和平区XX街123号", "1-23-4", "抵押", ""],
    ["2", "沈阳置业有限公司", "沈河区YY路45号", "2-15-1", "商铺", ""],
    ["3", "抚顺实业公司", "新抚区ZZ道67号", "3-8-2", "抵押", "已拍照"],
    ["4", "大连贸易有限公司", "中山区AA路89号", "5-3-7", "抵债", ""],
]
for r, row in enumerate(data):
    cx = tbl_x
    ry = tbl_y + Inches(0.45) + Inches(r * 0.5)
    bg = C_BG_LIGHT if r % 2 == 0 else C_WHITE
    for i, (val, cw) in enumerate(zip(row, col_widths)):
        cell = add_rect(slide, cx, ry, cw, Inches(0.5), bg, C_BORDER, Pt(0.5))
        add_text(slide, cx + Inches(0.05), ry, cw - Inches(0.1), Inches(0.5), val,
                 font_size=9, color=C_TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += cw

add_page_num(slide, 6)

# ============================================================
# Slide 6: 第四章 客户列表
# ============================================================
slide = add_slide(prs)
add_title_bar(slide, 4, "客户列表")

add_card(slide, Inches(0.4), Inches(1.0), Inches(6.0), Inches(5.5), "列表布局说明")
add_bullet_list(slide, Inches(0.6), Inches(1.6), Inches(5.6), Inches(4.8), [
    "表头区域：显示总客户数、已完成数、完成进度",
    "客户名：支持自动换行，完整显示长名称",
    "地址显示：地址概+地址精确组合，过长自动截断",
    "性质标签：住宅/商铺/抵押/抵债等，颜色区分",
    "拍照进度：显示各类型拍照数量（远X近X内X瑕X）",
    "备注按钮：点击添加/查看备注，已填备注显示●",
    "拍照按钮：进入拍照模式",
    "查看已拍：浏览该客户所有照片",
    "搜索功能：按客户名快速筛选",
], font_size=15)

# 右侧：列表模拟
lx = Inches(7.0)
ly = Inches(1.0)
lw = Inches(5.8)
lh = Inches(5.5)
add_phone_mockup(slide, lx, ly, lw, lh, "客户列表  总:4  已完成:2")
# 列表项
items = [
    ("成都投资集团", "和平区XX街123号", "远2近1内0瑕0", True),
    ("沈阳置业有限公司", "沈河区YY路45号2-15-1", "远1近2内0瑕0", True),
    ("抚顺实业公司", "新抚区ZZ道67号3-8-2", "远0近0内0瑕0", False),
    ("大连贸易有限公司", "中山区AA路89号5-3-7", "远0近0内0瑕0", False),
]
for i, (name, addr, progress, done) in enumerate(items):
    iy = ly + Inches(0.7) + Inches(i * 0.85)
    bg = C_CARD_BG if i % 2 == 0 else C_WHITE
    add_rect(slide, lx + Inches(0.15), iy, lw - Inches(0.3), Inches(0.75), bg, C_BORDER, Pt(0.5))
    name_color = C_SUCCESS if done else C_TEXT
    add_text(slide, lx + Inches(0.25), iy + Inches(0.05), Inches(1.8), Inches(0.3),
             name, font_size=9, bold=True, color=name_color)
    add_text(slide, lx + Inches(0.25), iy + Inches(0.35), Inches(2.5), Inches(0.3),
             addr, font_size=7, color=C_TEXT_DIM)
    add_text(slide, lx + Inches(2.2), iy + Inches(0.05), Inches(1.2), Inches(0.3),
             progress, font_size=8, color=C_SUCCESS if done else C_TEXT_DIM)
    btn_color = C_SUCCESS if done else C_WARNING
    add_rect(slide, lx + Inches(3.5), iy + Inches(0.1), Inches(0.9), Inches(0.28), btn_color)
    add_text(slide, lx + Inches(3.5), iy + Inches(0.1), Inches(0.9), Inches(0.28),
             "拍照" if not done else "已拍", font_size=8, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, lx + Inches(4.5), iy + Inches(0.1), Inches(0.8), Inches(0.28), C_ACCENT_D)
    add_text(slide, lx + Inches(4.5), iy + Inches(0.1), Inches(0.8), Inches(0.28),
             "查看", font_size=8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 底部日报表按钮
add_rect(slide, lx + Inches(0.15), ly + lh - Inches(0.7), lw - Inches(0.3), Inches(0.5),
         C_SUCCESS, radius=0.1)
add_text(slide, lx + Inches(0.15), ly + lh - Inches(0.7), lw - Inches(0.3), Inches(0.5),
         "AI 一键生成日报表", font_size=11, bold=True, color=C_WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_page_num(slide, 7)

# ============================================================
# Slide 7: 第五章 拍照功能
# ============================================================
slide = add_slide(prs)
add_title_bar(slide, 5, "拍照功能")

add_card(slide, Inches(0.4), Inches(1.0), Inches(6.0), Inches(5.5), "拍照流程")
add_bullet_list(slide, Inches(0.6), Inches(1.6), Inches(5.6), Inches(4.8), [
    "选择拍照类型：远景、近景、内部、瑕疵、其他",
    "手动快门连续拍摄：每按一次快门拍一张",
    "拍完点击对勾保存按钮确认",
    "照片自动添加水印（日期、地址、GPS）",
    "自动保存到应用目录并同步到系统相册",
    "文件名自动按规则生成（见第九章）",
    "支持返回后继续拍摄同一类型",
    "拍照进度实时更新（远X近X内X瑕X）",
], font_size=15)

# 右侧：拍照界面模拟
px = Inches(7.0)
py = Inches(1.0)
pw = Inches(5.8)
ph = Inches(5.5)
add_phone_mockup(slide, px, py, pw, ph, "成都投资集团 - 拍照")
# 类型选择
add_text(slide, px + Inches(0.2), py + Inches(0.65), pw - Inches(0.4), Inches(0.3),
         "请选择拍照类型：", font_size=10, color=C_TEXT)
types = ["远景", "近景", "内部", "瑕疵", "其他"]
for i, t in enumerate(types):
    tx = px + Inches(0.2) + Inches(i * 1.08)
    bg = C_ACCENT if i == 0 else C_BG_LIGHT
    tc = C_WHITE if i == 0 else C_TEXT
    add_rect(slide, tx, py + Inches(1.0), Inches(1.0), Inches(0.4), bg, C_BORDER, Pt(0.5))
    add_text(slide, tx, py + Inches(1.0), Inches(1.0), Inches(0.4),
             t, font_size=9, bold=True, color=tc, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 取景框
add_rect(slide, px + Inches(0.2), py + Inches(1.6), pw - Inches(0.4), Inches(2.8),
         C_BG_LIGHT, C_BORDER, Pt(1))
add_text(slide, px + Inches(0.2), py + Inches(2.5), pw - Inches(0.4), Inches(0.5),
         "[ 相机取景框 ]", font_size=12, color=C_TEXT_DIM, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 水印预览
wm_y = py + Inches(4.5)
add_rect(slide, px + Inches(0.2), wm_y, pw - Inches(0.4), Inches(0.4), C_BG_LIGHT, C_BORDER, Pt(0.5))
add_text(slide, px + Inches(0.3), wm_y, pw - Inches(0.6), Inches(0.4),
         "水印：2026年6月30日-和平区XX街123号-经纬度", font_size=8, color=C_TEXT_DIM, anchor=MSO_ANCHOR.MIDDLE)
# 状态
add_text(slide, px + Inches(0.2), py + Inches(5.0), pw - Inches(0.4), Inches(0.3),
         "当前：远景 02", font_size=10, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

add_page_num(slide, 8)

# ============================================================
# Slide 8: 第六章 照片查看与管理
# ============================================================
slide = add_slide(prs)
add_title_bar(slide, 6, "照片查看与管理")

add_card(slide, Inches(0.4), Inches(1.0), Inches(6.0), Inches(5.5), "照片管理功能")
add_bullet_list(slide, Inches(0.6), Inches(1.6), Inches(5.6), Inches(4.8), [
    "查看已拍照片：按类型分类展示（远景/近景/内部/瑕疵/其他）",
    "照片以缩略图网格排列，点击放大查看",
    "删除单张照片：长按或点击删除按钮",
    "删除全部照片：需二次确认，防止误删",
    "照片自动同步到系统相册（DCIM/Camera）",
    "可通过微信、文件管理器直接分享",
    "照片保存路径：应用私有目录 + 系统相册双备份",
], font_size=15)

# 右侧：照片查看模拟
gx = Inches(7.0)
gy = Inches(1.0)
gw = Inches(5.8)
gh = Inches(5.5)
add_phone_mockup(slide, gx, gy, gw, gh, "成都投资集团 - 已拍照片")
# 类型标签
tabs = [("远景(3)", True), ("近景(2)", False), ("内部(0)", False), ("瑕疵(1)", False)]
for i, (t, active) in enumerate(tabs):
    tx = gx + Inches(0.15) + Inches(i * 1.38)
    bg = C_ACCENT if active else C_BG_LIGHT
    tc = C_WHITE if active else C_TEXT_DIM
    add_rect(slide, tx, gy + Inches(0.65), Inches(1.3), Inches(0.35), bg, C_BORDER, Pt(0.5))
    add_text(slide, tx, gy + Inches(0.65), Inches(1.3), Inches(0.35),
             t, font_size=8, bold=True, color=tc, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 照片网格
for r in range(2):
    for c in range(3):
        idx = r * 3 + c
        if idx >= 3:
            break
        ph_x = gx + Inches(0.2) + Inches(c * 1.75)
        ph_y = gy + Inches(1.2) + Inches(r * 1.6)
        add_rect(slide, ph_x, ph_y, Inches(1.6), Inches(1.4), C_BG_LIGHT, C_BORDER, Pt(1))
        add_text(slide, ph_x, ph_y + Inches(0.5), Inches(1.6), Inches(0.4),
                 "照片%02d" % (idx + 1), font_size=9, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)
# 删除全部按钮
add_rect(slide, gx + Inches(0.2), gy + gh - Inches(0.8), gw - Inches(0.4), Inches(0.45),
         C_DANGER, radius=0.1)
add_text(slide, gx + Inches(0.2), gy + gh - Inches(0.8), gw - Inches(0.4), Inches(0.45),
         "删除全部（需二次确认）", font_size=10, bold=True, color=C_WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_page_num(slide, 9)

# ============================================================
# Slide 9: 第七章 设置页面
# ============================================================
slide = add_slide(prs)
add_title_bar(slide, 7, "设置页面")

add_card(slide, Inches(0.4), Inches(1.0), Inches(6.0), Inches(5.5), "设置项说明")
add_bullet_list(slide, Inches(0.6), Inches(1.6), Inches(5.6), Inches(4.8), [
    "命名规则：4段下拉选择器，组合照片文件名（见第九章）",
    "水印内容：3段内容选择 + 字号大小 + 水印位置",
    "水印字号：大(80pt) / 中(56pt) / 小(36pt)",
    "水印位置：左上 / 右上 / 左下 / 右下",
    "AI设置：API地址、API Key、模型ID",
    "保存设置后立即生效",
], font_size=15)

# 右侧：设置页模拟
sx = Inches(7.0)
sy = Inches(1.0)
sw = Inches(5.8)
sh = Inches(5.5)
add_phone_mockup(slide, sx, sy, sw, sh, "设置")
# 命名规则
add_text(slide, sx + Inches(0.2), sy + Inches(0.7), Inches(2.0), Inches(0.35),
         "命名规则：", font_size=10, bold=True, color=C_TEXT, anchor=MSO_ANCHOR.MIDDLE)
for i in range(4):
    dx = sx + Inches(2.2) + Inches(i * 0.85)
    add_rect(slide, dx, sy + Inches(0.72), Inches(0.75), Inches(0.3), C_BG_LIGHT, C_BORDER, Pt(0.5))
    add_text(slide, dx, sy + Inches(0.72), Inches(0.75), Inches(0.3),
             "段%d" % (i + 1), font_size=8, color=C_TEXT_DIM, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 水印内容
add_text(slide, sx + Inches(0.2), sy + Inches(1.3), Inches(2.0), Inches(0.35),
         "水印内容：", font_size=10, bold=True, color=C_TEXT, anchor=MSO_ANCHOR.MIDDLE)
for i in range(3):
    dx = sx + Inches(2.2) + Inches(i * 1.1)
    add_rect(slide, dx, sy + Inches(1.32), Inches(1.0), Inches(0.3), C_BG_LIGHT, C_BORDER, Pt(0.5))
    add_text(slide, dx, sy + Inches(1.32), Inches(1.0), Inches(0.3),
             "水%d" % (i + 1), font_size=8, color=C_TEXT_DIM, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 水印字号
add_text(slide, sx + Inches(0.2), sy + Inches(1.9), Inches(2.0), Inches(0.35),
         "水印字号：", font_size=10, bold=True, color=C_TEXT, anchor=MSO_ANCHOR.MIDDLE)
for i, (label, active) in enumerate([("大", True), ("中", False), ("小", False)]):
    dx = sx + Inches(2.2) + Inches(i * 0.9)
    bg = C_ACCENT if active else C_BG_LIGHT
    tc = C_WHITE if active else C_TEXT
    add_rect(slide, dx, sy + Inches(1.92), Inches(0.8), Inches(0.3), bg, C_BORDER, Pt(0.5))
    add_text(slide, dx, sy + Inches(1.92), Inches(0.8), Inches(0.3),
             label, font_size=9, bold=True, color=tc, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 水印位置
add_text(slide, sx + Inches(0.2), sy + Inches(2.5), Inches(2.0), Inches(0.35),
         "水印位置：", font_size=10, bold=True, color=C_TEXT, anchor=MSO_ANCHOR.MIDDLE)
for i, label in enumerate(["左上", "右上", "左下", "右下"]):
    dx = sx + Inches(2.2) + Inches(i * 0.85)
    bg = C_ACCENT if i == 0 else C_BG_LIGHT
    tc = C_WHITE if i == 0 else C_TEXT
    add_rect(slide, dx, sy + Inches(2.52), Inches(0.75), Inches(0.3), bg, C_BORDER, Pt(0.5))
    add_text(slide, dx, sy + Inches(2.52), Inches(0.75), Inches(0.3),
             label, font_size=8, bold=True, color=tc, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# AI设置
add_rect(slide, sx + Inches(0.15), sy + Inches(3.2), sw - Inches(0.3), Inches(0.35), C_BG_LIGHT)
add_text(slide, sx + Inches(0.25), sy + Inches(3.2), sw - Inches(0.5), Inches(0.35),
         "AI 助手设置", font_size=10, bold=True, color=C_ACCENT_D, anchor=MSO_ANCHOR.MIDDLE)
add_text(slide, sx + Inches(0.25), sy + Inches(3.6), Inches(2.0), Inches(0.3),
         "API地址", font_size=9, color=C_TEXT_DIM, anchor=MSO_ANCHOR.MIDDLE)
add_rect(slide, sx + Inches(2.2), sy + Inches(3.62), Inches(3.3), Inches(0.3), C_BG_LIGHT, C_BORDER, Pt(0.5))
add_text(slide, sx + Inches(2.3), sy + Inches(3.62), Inches(3.0), Inches(0.3),
         "https://api.deepseek.com/v1", font_size=7, color=C_TEXT_DIM, anchor=MSO_ANCHOR.MIDDLE)
# 保存按钮
add_rect(slide, sx + Inches(1.0), sy + sh - Inches(0.8), sw - Inches(2.0), Inches(0.5),
         C_ACCENT, radius=0.1)
add_text(slide, sx + Inches(1.0), sy + sh - Inches(0.8), sw - Inches(2.0), Inches(0.5),
         "保 存 设 置", font_size=12, bold=True, color=C_WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_page_num(slide, 10)

# ============================================================
# Slide 10: 第八章 水印说明
# ============================================================
slide = add_slide(prs)
add_title_bar(slide, 8, "水印说明")

add_card(slide, Inches(0.4), Inches(1.0), Inches(6.0), Inches(5.5), "水印配置详解")
add_bullet_list(slide, Inches(0.6), Inches(1.6), Inches(5.6), Inches(4.8), [
    "每段水印可从以下选项中选择：",
    "  - 拍摄日期（仅年月日，如：2026年6月30日）",
    "  - 客户名称",
    "  - 地址（概+精确拼接）",
    "  - 经纬度（GPS定位信息）",
    "  - 空值（该段不显示）",
    "字号选择：大(80pt) / 中(56pt) / 小(36pt)",
    "位置选择：左上 / 右上 / 左下 / 右下",
    "v3.19起：水印区域及字体增大2倍，日期仅年月日",
], font_size=15)

# 右侧：水印位置示意图（矢量图形，替换损坏的生成图）
wm_frame_x = Inches(7.0)
wm_frame_y = Inches(1.0)
wm_frame_w = Inches(5.8)
wm_frame_h = Inches(4.0)
# 模拟照片框
add_rect(slide, wm_frame_x, wm_frame_y, wm_frame_w, wm_frame_h, C_BG_LIGHT, C_BORDER, Pt(1.5))
# 模拟照片内容（渐变色块）
add_rect(slide, wm_frame_x + Inches(0.3), wm_frame_y + Inches(0.3), Inches(2.0), Inches(1.5), RGBColor(0xD6, 0xEA, 0xFA))
add_rect(slide, wm_frame_x + Inches(2.5), wm_frame_y + Inches(0.3), Inches(3.0), Inches(1.5), RGBColor(0xE8, 0xF5, 0xE9))
add_rect(slide, wm_frame_x + Inches(0.3), wm_frame_y + Inches(2.0), Inches(5.2), Inches(1.5), RGBColor(0xFE, 0xF3, 0xE2))
# 四个水印位置标注
positions = [
    ("左上", Inches(0.4), Inches(0.4), C_ACCENT),
    ("右上", wm_frame_w - Inches(1.4), Inches(0.4), C_SUCCESS),
    ("左下", Inches(0.4), wm_frame_h - Inches(0.6), C_WARNING),
    ("右下", wm_frame_w - Inches(1.4), wm_frame_h - Inches(0.6), C_DANGER),
]
for label, dx, dy, color in positions:
    px = wm_frame_x + dx
    py = wm_frame_y + dy
    add_rect(slide, px, py, Inches(1.0), Inches(0.35), color, radius=0.1)
    add_text(slide, px, py, Inches(1.0), Inches(0.35),
             label, font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 水印文字示例
add_text(slide, wm_frame_x + Inches(0.3), wm_frame_y + Inches(2.0), Inches(5.2), Inches(1.5),
         "水印示例：\n2026年6月30日\n成都投资集团\n和平区XX街123号\nN41.78° E123.45°",
         font_size=12, bold=True, color=RGBColor(0x33, 0x33, 0x33), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 水印位置标注
add_card(slide, Inches(7.0), Inches(5.2), Inches(5.8), Inches(1.3), "水印位置示意")
add_text(slide, Inches(7.2), Inches(5.6), Inches(5.4), Inches(0.8),
         "左上 ←→ 右上\n左下 ←→ 右下\n（点击设置页选择水印位置）", font_size=11, color=C_TEXT_DIM)

add_page_num(slide, 11)

# ============================================================
# Slide 11: 第九章 文件命名规则
# ============================================================
slide = add_slide(prs)
add_title_bar(slide, 9, "文件命名规则")

add_card(slide, Inches(0.4), Inches(1.0), Inches(12.5), Inches(1.8), "命名规则")
add_text(slide, Inches(0.6), Inches(1.5), Inches(12.0), Inches(1.2),
         ["文件名由4个可配置字段段组合而成，各段用连字符「-」连接",
          "每段可选：拍摄日期 / 客户名 / 地址+时间 / 空值（自动省略）",
          "文件名末尾自动追加：类型（远景/近景/内部/瑕疵/其他）+ 编号（01/02/03…）"],
         font_size=15, color=C_TEXT)

# 命名结构图
diag_y = Inches(3.2)
segments = [
    ("段1\n拍摄日期", C_ACCENT),
    ("段2\n客户名", C_ACCENT_D),
    ("段3\n地址+时间", C_SUCCESS),
    ("段4\n空值", C_TEXT_DIM),
]
seg_w = Inches(2.2)
seg_gap = Inches(0.25)
start_x = Inches(0.8)
for i, (label, color) in enumerate(segments):
    sx = start_x + i * (seg_w + seg_gap)
    add_rect(slide, sx, diag_y, seg_w, Inches(0.8), color, radius=0.08)
    add_text(slide, sx, diag_y, seg_w, Inches(0.8),
             label, font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 3:
        add_text(slide, sx + seg_w, diag_y, seg_gap, Inches(0.8),
                 "-", font_size=20, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 类型+编号
add_text(slide, start_x + 4 * (seg_w + seg_gap), diag_y, Inches(1.5), Inches(0.8),
         "+ 类型 + 编号", font_size=11, bold=True, color=C_TEXT, anchor=MSO_ANCHOR.MIDDLE)

# 示例
add_card(slide, Inches(0.4), Inches(4.5), Inches(12.5), Inches(1.5), "示例")
add_rect(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(0.6), C_BG_LIGHT, C_BORDER, Pt(0.5))
add_text(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(0.6),
         "20260630-成都投资集团-和平区XX街123号1430-远景-01.jpg",
         font_size=14, bold=True, color=C_ACCENT_D, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_page_num(slide, 12)

# ============================================================
# Slide 12: 第十章 日报表生成功能（新增）
# ============================================================
slide = add_slide(prs)
add_title_bar(slide, 10, "日报表生成功能")

# 流程图
steps = [
    ("1", "收集数据", "仅对有外访照片的\n客户汇总名称\n地址、类型、备注\n拍照数量"),
    ("2", "DeepSeek生成", "调用DeepSeek\n为已外访客户撰写\n抵押物情况\n现状描述\n风险备注"),
    ("3", "填充模板+汇总", "填入内置日报表模板\n末尾追加汇总说明\n（基于XX文件生成\n共计XX户/外访XX户）"),
    ("4", "自动命名保存", "默认文件名：\n抵押物、抵债资产\n现场勘查日报表\nYYYYMMDD.xlsx"),
]
step_w = Inches(2.8)
step_gap = Inches(0.25)
start_x = Inches(0.6)
step_y = Inches(1.2)
for i, (num, title, desc) in enumerate(steps):
    sx = start_x + i * (step_w + step_gap)
    # 卡片
    add_rect(slide, sx, step_y, step_w, Inches(3.0), C_CARD_BG, C_BORDER, Pt(1), radius=0.05)
    # 编号圆
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, sx + step_w/2 - Inches(0.3), step_y + Inches(0.15),
                                   Inches(0.6), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = C_ACCENT
    badge.line.fill.background()
    add_text_in_shape(badge, num, font_size=18, bold=True, color=C_WHITE)
    # 标题
    add_text(slide, sx, step_y + Inches(0.8), step_w, Inches(0.4),
             title, font_size=14, bold=True, color=C_ACCENT_D, align=PP_ALIGN.CENTER)
    # 描述
    add_text(slide, sx + Inches(0.1), step_y + Inches(1.3), step_w - Inches(0.2), Inches(1.5),
             desc, font_size=11, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)
    # 箭头
    if i < 3:
        ax = sx + step_w + Inches(0.02)
        add_text(slide, ax, step_y + Inches(1.2), step_gap, Inches(0.5),
                 "→", font_size=20, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 报表模拟（矢量图形，替换损坏的生成图）
rpt_x = Inches(0.6)
rpt_y = Inches(4.5)
rpt_w = Inches(4.5)
rpt_h = Inches(2.5)
add_rect(slide, rpt_x, rpt_y, rpt_w, rpt_h, C_WHITE, C_BORDER, Pt(1.5), radius=0.02)
# 表头
add_rect(slide, rpt_x, rpt_y, rpt_w, Inches(0.35), C_ACCENT)
add_text(slide, rpt_x, rpt_y, rpt_w, Inches(0.35),
         "抵押物、抵债资产现场勘查日报表", font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 表格行
rpt_headers = ["序号", "贷款人", "抵押物情况", "风险"]
rpt_col_w = [Inches(0.6), Inches(1.4), Inches(1.7), Inches(0.8)]
cx = rpt_x
for hdr, cw in zip(rpt_headers, rpt_col_w):
    add_rect(slide, cx, rpt_y + Inches(0.35), cw, Inches(0.3), C_BG_LIGHT, C_BORDER, Pt(0.5))
    add_text(slide, cx, rpt_y + Inches(0.35), cw, Inches(0.3),
             hdr, font_size=8, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += cw
# 数据行
rpt_data = [
    ["1", "成都投资集团", "办公楼正常", "无"],
    ["2", "沈阳置业有限", "商铺正常", "无"],
    ["3", "抚顺实业公司", "待核实", "关注"],
]
for r, row in enumerate(rpt_data):
    cx = rpt_x
    ry = rpt_y + Inches(0.65) + Inches(r * 0.35)
    for val, cw in zip(row, rpt_col_w):
        add_rect(slide, cx, ry, cw, Inches(0.35), C_WHITE if r % 2 == 0 else C_BG_LIGHT, C_BORDER, Pt(0.5))
        risk_color = C_DANGER if val == "关注" else C_TEXT
        add_text(slide, cx, ry, cw, Inches(0.35),
                 val, font_size=8, color=risk_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += cw
# 汇总说明行（合并整行，紧接数据末尾）
sum_y = rpt_y + Inches(1.75)
add_rect(slide, rpt_x, sum_y, rpt_w, Inches(0.3), RGBColor(0xE3, 0xF2, 0xFD), C_BORDER, Pt(0.5))
add_text(slide, rpt_x + Inches(0.05), sum_y, rpt_w - Inches(0.1), Inches(0.3),
         "本次基于客户清单.xlsx生成，共计4户/外访3户已生成报告",
         font_size=7, bold=True, color=C_ACCENT_D, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 签字栏
add_text(slide, rpt_x + Inches(0.1), rpt_y + rpt_h - Inches(0.4), rpt_w - Inches(0.2), Inches(0.3),
         "填报人：____  日期：2026-06-30", font_size=8, color=C_TEXT_DIM)

# 注意事项
add_card(slide, Inches(5.5), Inches(4.5), Inches(7.2), Inches(2.5), "注意事项")
add_bullet_list(slide, Inches(5.7), Inches(5.0), Inches(6.8), Inches(1.8), [
    "仅对有外访照片的客户生成报告行（未外访客户不计入）",
    "末尾自动追加汇总说明：基于XX文件生成，共计XX户/外访XX户",
    "自动命名「抵押物、抵债资产现场勘查日报表YYYYMMDD.xlsx」",
    "内置DeepSeek(deepseek-v4-flash)模型，无需额外配置",
    "报告中完整列出所有人名和地址，不使用「等」字省略",
    "点击后按钮显示「正在生成中…」避免重复点击",
    "AI生成内容仅供参考，建议人工审核后使用",
], font_size=11)

add_page_num(slide, 13)

# ============================================================
# Slide 13: 第十一章 AI助手功能
# ============================================================
slide = add_slide(prs)
add_title_bar(slide, 11, "AI助手功能")

add_card(slide, Inches(0.4), Inches(1.0), Inches(6.0), Inches(5.5), "AI助手使用说明")
add_bullet_list(slide, Inches(0.6), Inches(1.6), Inches(5.6), Inches(4.8), [
    "点击主界面「AI助手」按钮进入聊天界面",
    "可直接询问拍摄进度：",
    "  - 今天拍了多少照片？",
    "  - 某某公司拍了没有？",
    "  - 远景拍了多少张？",
    "内置DeepSeek(deepseek-v4-flash)模型，无需额外配置即可使用",
    "也可在设置页面自定义API地址和模型",
    "AI助手帮助快速查询盘点进度，无需手动统计",
], font_size=15)

# 右侧：聊天界面模拟
cx = Inches(7.0)
cy = Inches(1.0)
cw = Inches(5.8)
ch = Inches(5.5)
add_phone_mockup(slide, cx, cy, cw, ch, "AI 拍摄助手")
# 对话气泡
msgs = [
    ("user", "今天拍了多少照片？"),
    ("assistant", "今天共拍摄了 8 张照片：\n成都投资集团：远2近1\n沈阳置业：远1近2\n抚顺实业：远0近0\n大连贸易：远0近0"),
    ("user", "远景拍了多少张？"),
    ("assistant", "远景共拍摄 3 张。"),
]
msg_y = cy + Inches(0.7)
for role, text in msgs:
    lines = text.count('\n') + 1
    h = Inches(0.3 + lines * 0.22)
    if role == "user":
        bx = cx + cw - Inches(3.5) - Inches(0.2)
        bg = C_ACCENT
        tc = C_WHITE
        align = PP_ALIGN.RIGHT
    else:
        bx = cx + Inches(0.2)
        bg = C_BG_LIGHT
        tc = C_TEXT
        align = PP_ALIGN.LEFT
    add_rect(slide, bx, msg_y, Inches(3.5), h, bg, C_BORDER, Pt(0.5), radius=0.1)
    add_text(slide, bx + Inches(0.1), msg_y + Inches(0.05), Inches(3.3), h - Inches(0.1),
             text, font_size=8, color=tc, align=align)
    msg_y += h + Inches(0.15)

add_page_num(slide, 14)

# ============================================================
# Slide 14: 第十二章 日志管理（v3.22.2 新增）
# ============================================================
slide = add_slide(prs)
add_title_bar(slide, 12, "日志管理")

# 左侧：日志功能说明
add_card(slide, Inches(0.4), Inches(1.0), Inches(6.0), Inches(5.5), "日志功能说明")
add_bullet_list(slide, Inches(0.6), Inches(1.6), Inches(5.6), Inches(4.8), [
    "主界面底部「日志」按钮点击后弹出二级管理弹窗",
    "弹窗含 4 个功能按钮 + 当前开关状态显示",
    "开关：开启/关闭全量 app 日志记录（默认关闭）",
    "查看：打开三级弹窗显示完整日志内容（可滚动）",
    "复制：将日志全文复制到剪贴板，便于反馈问题",
    "清空：删除日志文件（不可恢复，需谨慎）",
    "v3.22.2 起：UI 交互同步入全量日志（含查看已拍、搜索、保存备注等）",
    "日志文件路径：应用私有目录 app_debug.log",
    "开关状态持久化，重启后保持",
], font_size=14)

# 右侧：日志管理弹窗模拟
lgx = Inches(7.0)
lgy = Inches(1.0)
lgw = Inches(5.8)
lgh = Inches(5.5)
add_phone_mockup(slide, lgx, lgy, lgw, lgh, "日志管理")
# 标题
add_text(slide, lgx + Inches(0.2), lgy + Inches(0.65), lgw - Inches(0.4), Inches(0.4),
         "日志管理", font_size=14, bold=True, color=C_ACCENT_D, anchor=MSO_ANCHOR.MIDDLE)
# 状态
add_text(slide, lgx + Inches(0.2), lgy + Inches(1.1), lgw - Inches(0.4), Inches(0.35),
         "当前状态: 已关闭", font_size=11, color=C_TEXT_DIM, anchor=MSO_ANCHOR.MIDDLE)
# 4 按钮 2x2 网格
btn_w = Inches(2.5)
btn_h = Inches(0.7)
btn_gap = Inches(0.2)
btn_start_x = lgx + Inches(0.3)
btn_start_y = lgy + Inches(1.6)
btns = [
    ("开启日志记录", C_SUCCESS),
    ("查看日志", C_ACCENT),
    ("复制日志", C_ACCENT_D),
    ("清空日志", C_WARNING),
]
for i, (label, color) in enumerate(btns):
    col = i % 2
    row = i // 2
    bx = btn_start_x + col * (btn_w + btn_gap)
    by = btn_start_y + row * (btn_h + btn_gap)
    add_rect(slide, bx, by, btn_w, btn_h, color, radius=0.1)
    add_text(slide, bx, by, btn_w, btn_h,
             label, font_size=12, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 关闭按钮
add_rect(slide, lgx + Inches(0.3), lgy + lgh - Inches(1.0), lgw - Inches(0.6), Inches(0.5),
         C_BORDER, radius=0.1)
add_text(slide, lgx + Inches(0.3), lgy + lgh - Inches(1.0), lgw - Inches(0.6), Inches(0.5),
         "关闭", font_size=12, color=C_TEXT_DIM,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 说明
add_text(slide, lgx + Inches(0.2), lgy + lgh - Inches(0.4), lgw - Inches(0.4), Inches(0.3),
         "提示：开启日志后所有操作将被记录，便于问题排查", font_size=9, color=C_TEXT_DIM,
         align=PP_ALIGN.CENTER)

add_page_num(slide, 15)

# ============================================================
# Slide 15: 第十三章 常见问题 + 末页
# ============================================================
slide = add_slide(prs)
add_title_bar(slide, 13, "常见问题")

faqs = [
    ("Q1: 点击拍照后相机没有启动？",
     "请确认已授予相机权限（手机设置→应用管理→本应用→权限管理）。也可尝试在拍照类型弹窗中选择「其他」类型启动。"),
    ("Q2: GPS定位不准或无法获取？",
     "请确保手机已开启GPS定位功能，并在开阔地带使用。水印中定位信息需等待GPS定位成功后显示，首次定位可能需要30秒。"),
    ("Q3: 照片保存在哪里？",
     "照片保存在应用私有目录下，同时自动复制到系统相册（DCIM/Camera）。可在手机相册、微信发送图片中直接查看。"),
    ("Q4: 如何一键生成勘查日报表？",
     "点击主界面底部「AI一键生成日报表」按钮，AI将根据备注和Excel数据自动填写内置模板，生成后可选择保存位置。"),
    ("Q5: 删除照片能否恢复？",
     "不能。删除前会弹出二次确认对话框，确认后永久删除且无法恢复，请谨慎操作。"),
]

faq_y = Inches(1.0)
for i, (q, a) in enumerate(faqs):
    # 问题
    add_rect(slide, Inches(0.4), faq_y, Inches(12.5), Inches(0.35), C_BG_LIGHT, C_BORDER, Pt(0.5))
    add_text(slide, Inches(0.5), faq_y, Inches(12.3), Inches(0.35),
             q, font_size=12, bold=True, color=C_ACCENT_D, anchor=MSO_ANCHOR.MIDDLE)
    # 答案
    add_text(slide, Inches(0.6), faq_y + Inches(0.38), Inches(12.2), Inches(0.7),
             a, font_size=11, color=C_TEXT_DIM)
    faq_y += Inches(1.05)

# 底部联系信息
add_rect(slide, Inches(0.4), faq_y + Inches(0.1), Inches(12.5), Inches(0.55), C_ACCENT, radius=0.05)
add_text(slide, Inches(0.5), faq_y + Inches(0.1), Inches(12.3), Inches(0.55),
         "问题咨询电话：15940454123    |    资产盘点专项拍照工具 v3.22.24",
         font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_page_num(slide, 16)

# ============================================================
# 保存
# ============================================================
prs.save(DST)
print(f"保存成功: {DST}")
print(f"文件大小: {os.path.getsize(DST)} bytes")
print(f"总页数: {len(prs.slides)}")
