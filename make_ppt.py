from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG_COLOR = RGBColor(0x1C, 0x1C, 0x24)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x2D, 0x2D, 0x3A)
ACCENT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
SCREEN_BG = RGBColor(0x25, 0x25, 0x30)
SCREEN_BORDER = RGBColor(0x55, 0x55, 0x66)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = '微软雅黑'
    p.alignment = align
    return txBox


def add_paragraph(tf, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = '微软雅黑'
    p.alignment = align
    p.space_before = space_before
    return p


def add_phone_screen(slide, left, top, width, height, title_text=""):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SCREEN_BG
    shape.line.color.rgb = SCREEN_BORDER
    shape.line.width = Pt(2)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    
    p = tf.paragraphs[0]
    p.text = title_text if title_text else "手机屏幕模拟"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER
    
    return shape


def add_screen_element(slide, left, top, width, height, text, bg_color=DARK_GRAY, text_color=WHITE, font_size=9):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0x44, 0x44, 0x55)
    shape.line.width = Pt(1)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return shape


def add_page_number(slide, page_num, total_pages):
    add_textbox(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
                f"{page_num} / {total_pages}", font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


def add_title_bar(slide, title_text):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                title_text, font_size=32, bold=True, color=WHITE)
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    total_slides = 13
    
    blank_layout = prs.slide_layouts[6]
    
    # ==================== 第1页：封面 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    for i in range(5):
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5 + i * 0.15), Inches(2.8 - i * 0.08),
            Inches(4 - i * 0.3), Inches(0.04)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = ACCENT_BLUE if i == 0 else RGBColor(0x3A + i*5, 0x70 + i*5, 0xB0 + i*5)
        accent_line.line.fill.background()
    
    add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                "资产盘点专项拍照工具", font_size=44, bold=True, color=WHITE)
    
    add_textbox(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.7),
                "使用说明书 v3.6.0", font_size=28, color=ACCENT_BLUE)
    
    add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                "2026年6月", font_size=18, color=LIGHT_GRAY)
    
    add_page_number(slide, 1, total_slides)
    
    # ==================== 第2页：目录 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_bar(slide, "目录")
    
    chapters = [
        "第一章  应用概述",
        "第二章  首次使用",
        "第三章  Excel导入",
        "第四章  客户列表",
        "第五章  拍照功能",
        "第六章  照片查看与管理",
        "第七章  设置页面",
        "第八章  水印说明",
        "第九章  文件命名规则",
        "第十章  常见问题",
    ]
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, chapter in enumerate(chapters):
        col = i // 5
        row = i % 5
        x = Inches(1 + col * 5.5)
        y = Inches(1.8 + row * 0.9)
        
        num_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.5), Inches(0.5))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = ACCENT_BLUE
        num_box.line.fill.background()
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]
        np.text = str(i + 1)
        np.font.size = Pt(16)
        np.font.color.rgb = WHITE
        np.font.bold = True
        np.font.name = '微软雅黑'
        np.alignment = PP_ALIGN.CENTER
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        add_textbox(slide, x + Inches(0.7), y + Inches(0.08), Inches(4.5), Inches(0.5),
                    chapter, font_size=20, color=WHITE)
    
    add_page_number(slide, 2, total_slides)
    
    # ==================== 第3页：第一章 应用概述 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_bar(slide, "第一章  应用概述")
    
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(6.5), Inches(0.5),
                "功能介绍", font_size=24, bold=True, color=ACCENT_BLUE)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(6.5), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "资产盘点专项拍照工具是一款专为信贷资产盘点工作设计的Android拍照应用，主要功能包括："
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.name = '微软雅黑'
    
    features = [
        "Excel批量导入客户地址信息",
        "按客户分类管理资产照片",
        "支持远景/近景/内部/瑕疵四种拍照类型",
        "自动添加水印（时间、地址、GPS定位）",
        "自定义文件命名规则",
        "手动快门连续拍摄，按返回键结束",
        "删除照片二次确认，防止误操作",
    ]
    for feat in features:
        add_paragraph(tf, "  - " + feat, font_size=17, color=LIGHT_GRAY, space_before=Pt(8))
    
    add_textbox(slide, Inches(0.5), Inches(4.8), Inches(6.5), Inches(0.5),
                "适用场景", font_size=24, bold=True, color=ACCENT_BLUE)
    
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(6.5), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "适用于银行、小贷公司、资产管理公司等金融机构的贷前调查、贷中检查、贷后催收、资产保全等需要现场拍照取证的业务场景。"
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.font.name = '微软雅黑'
    
    screen_left = Inches(7.5)
    screen_top = Inches(1.6)
    screen_w = Inches(4.8)
    screen_h = Inches(5.2)
    phone = add_phone_screen(slide, screen_left, screen_top, screen_w, screen_h, "应用主界面")
    
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(0.6), screen_w - Inches(0.6), Inches(0.7),
                       "资产盘点拍照工具", bg_color=ACCENT_BLUE, font_size=14)
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(1.6), screen_w - Inches(0.6), Inches(0.6),
                       "[导入Excel]  按钮", font_size=11)
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(2.4), screen_w - Inches(0.6), Inches(0.6),
                       "[设置]     按钮", font_size=11)
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(3.3), screen_w - Inches(0.6), Inches(1.5),
                       "客户列表区域\n显示所有导入的客户\n点击进入拍照", font_size=10)
    
    add_page_number(slide, 3, total_slides)
    
    # ==================== 第4页：第二章 首次使用 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_bar(slide, "第二章  首次使用")
    
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(6.5), Inches(0.5),
                "欢迎页介绍", font_size=24, bold=True, color=ACCENT_BLUE)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(6.5), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "首次启动应用时，会显示欢迎引导页面，简要介绍应用的主要功能和使用流程。用户可滑动查看，点击「开始使用」按钮进入主界面。"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.name = '微软雅黑'
    
    add_textbox(slide, Inches(0.5), Inches(4.0), Inches(6.5), Inches(0.5),
                "权限说明", font_size=24, bold=True, color=ACCENT_BLUE)
    
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(6.5), Inches(2.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "应用运行需要以下权限，请在首次使用时全部允许："
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.font.name = '微软雅黑'
    
    perms = [
        "相机权限：用于拍摄资产照片",
        "存储权限：用于保存照片到手机",
        "位置权限：用于获取GPS定位信息（水印用）",
        "文件读取权限：用于读取Excel文件",
    ]
    for perm in perms:
        add_paragraph(tf2, "  - " + perm, font_size=17, color=LIGHT_GRAY, space_before=Pt(8))
    
    screen_left = Inches(7.5)
    screen_top = Inches(1.6)
    screen_w = Inches(4.8)
    screen_h = Inches(5.2)
    phone = add_phone_screen(slide, screen_left, screen_top, screen_w, screen_h, "欢迎页 / 权限申请")
    
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(0.6), screen_w - Inches(0.6), Inches(1.0),
                       "欢迎使用\n资产盘点拍照工具", bg_color=DARK_GRAY, font_size=12)
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(1.9), screen_w - Inches(0.6), Inches(1.5),
                       "引导图区域\n展示功能介绍", font_size=10)
    add_screen_element(slide, screen_left + Inches(0.8), screen_top + Inches(3.7), screen_w - Inches(1.6), Inches(0.6),
                       "[ 开始使用 ]", bg_color=ACCENT_BLUE, font_size=12)
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(4.6), screen_w - Inches(0.6), Inches(0.5),
                       "允许相机权限？", font_size=10)
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(5.2), screen_w - Inches(1.6) / 2 - Inches(0.15), Inches(0.4),
                       "拒绝", font_size=10)
    add_screen_element(slide, screen_left + screen_w/2 + Inches(0.15), screen_top + Inches(5.2), screen_w/2 - Inches(0.45), Inches(0.4),
                       "允许", bg_color=ACCENT_BLUE, font_size=10)
    
    add_page_number(slide, 4, total_slides)
    
    # ==================== 第5页：第三章 Excel导入 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_bar(slide, "第三章  Excel导入")
    
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(6.5), Inches(0.5),
                "Excel格式要求", font_size=24, bold=True, color=ACCENT_BLUE)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(6.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Excel文件必须包含以下列（第一行为表头）："
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.name = '微软雅黑'
    
    cols = [
        "A列：客户名",
        "B列：地址概（如：沈阳市和平区）",
        "C列：地址精确（详细门牌号）",
        "D列：性质（如：住宅/商铺/抵押）",
    ]
    for col in cols:
        add_paragraph(tf, "  - " + col, font_size=17, color=LIGHT_GRAY, space_before=Pt(6))
    
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(6.5), Inches(0.5),
                "操作步骤", font_size=24, bold=True, color=ACCENT_BLUE)
    
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(6.5), Inches(2.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    steps = [
        "1. 点击主界面「导入Excel」按钮",
        "2. 在文件浏览器中选择准备好的.xlsx文件",
        "3. 应用自动解析文件内容",
        "4. 导入成功后显示客户数量提示",
        "5. 客户列表自动刷新显示新导入的数据",
    ]
    p2 = tf2.paragraphs[0]
    p2.text = steps[0]
    p2.font.size = Pt(17)
    p2.font.color.rgb = WHITE
    p2.font.name = '微软雅黑'
    for step in steps[1:]:
        add_paragraph(tf2, step, font_size=17, color=LIGHT_GRAY, space_before=Pt(8))
    
    screen_left = Inches(7.5)
    screen_top = Inches(1.6)
    screen_w = Inches(4.8)
    screen_h = Inches(5.2)
    phone = add_phone_screen(slide, screen_left, screen_top, screen_w, screen_h, "Excel导入界面")
    
    add_screen_element(slide, screen_left + Inches(0.2), screen_top + Inches(0.5), screen_w - Inches(0.4), Inches(0.4),
                       " <- 返回   导入Excel", bg_color=ACCENT_BLUE, font_size=11)
    
    excel_x = screen_left + Inches(0.3)
    excel_y = screen_top + Inches(1.1)
    col_w = (screen_w - Inches(0.6)) / 4
    headers = ["A客户名", "B地址概", "C地址精确", "D性质"]
    for i, h in enumerate(headers):
        add_screen_element(slide, excel_x + col_w * i, excel_y, col_w, Inches(0.35),
                           h, bg_color=RGBColor(0x44,0x44,0x55), font_size=8)
    rows = [
        ["张三", "和平区", "XX街123号", "住宅"],
        ["李四", "沈河区", "YY路45号", "商铺"],
        ["王五", "铁西区", "ZZ道67号", "抵押"],
    ]
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            add_screen_element(slide, excel_x + col_w * c_idx, excel_y + Inches(0.35) * (r_idx + 1),
                               col_w, Inches(0.35), val, font_size=8)
    
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(3.2), screen_w - Inches(0.6), Inches(0.5),
                       "[ 选择文件 ]", bg_color=ACCENT_BLUE, font_size=12)
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(4.0), screen_w - Inches(0.6), Inches(0.6),
                       "导入成功！\n共导入 3 条客户数据", bg_color=RGBColor(0x2D,0x5A,0x2D), font_size=10)
    
    add_page_number(slide, 5, total_slides)
    
    # ==================== 第6页：第四章 客户列表 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_bar(slide, "第四章  客户列表")
    
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(6.5), Inches(0.5),
                "列表布局说明", font_size=24, bold=True, color=ACCENT_BLUE)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(6.5), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    items = [
        ("表头区域：", "显示总客户数、已完成数、完成进度百分比"),
        ("客户名显示：", "客户名支持自动换行，完整显示长名称"),
        ("地址显示：", "优先显示地址概+地址精确组合，地址过长自动截断"),
        ("性质标签：", "用不同颜色标签区分房产性质"),
        ("拍照按钮：", "点击进入该客户的拍照页面，选择类型后启动相机"),
        ("查看已拍：", "查看该客户已拍摄的照片，支持删除（需二次确认）"),
        ("进度统计：", "实时显示每个客户的拍照进度（远景/近景/内部/瑕疵）"),
    ]
    
    p = tf.paragraphs[0]
    p.text = items[0][0] + items[0][1]
    p.font.size = Pt(17)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = '微软雅黑'
    for title, desc in items[1:]:
        p = add_paragraph(tf, title + desc, font_size=17, color=LIGHT_GRAY, space_before=Pt(10))
        p.runs[0].font.bold = True
    
    screen_left = Inches(7.5)
    screen_top = Inches(1.5)
    screen_w = Inches(4.8)
    screen_h = Inches(5.4)
    phone = add_phone_screen(slide, screen_left, screen_top, screen_w, screen_h, "客户列表页")
    
    add_screen_element(slide, screen_left + Inches(0.2), screen_top + Inches(0.5), screen_w - Inches(0.4), Inches(0.45),
                       "客户列表  总:3  已完成:1", bg_color=ACCENT_BLUE, font_size=10)
    
    add_screen_element(slide, screen_left + Inches(0.2), screen_top + Inches(1.1), screen_w - Inches(0.4), Inches(0.3),
                       "客户名          地址                     进度", bg_color=RGBColor(0x44,0x44,0x55), font_size=8)
    
    list_y = screen_top + Inches(1.5)
    for i in range(3):
        bg = RGBColor(0x30,0x30,0x3D) if i % 2 == 0 else SCREEN_BG
        add_screen_element(slide, screen_left + Inches(0.2), list_y + Inches(0.95)*i, screen_w - Inches(0.4), Inches(0.85),
                           "", bg_color=bg, font_size=8)
        add_screen_element(slide, screen_left + Inches(0.3), list_y + Inches(0.95)*i + Inches(0.05), Inches(1.0), Inches(0.35),
                           ["张三", "李四", "王五"][i], bg_color=bg, text_color=WHITE, font_size=10)
        add_screen_element(slide, screen_left + Inches(1.35), list_y + Inches(0.95)*i + Inches(0.05), Inches(1.8), Inches(0.35),
                           ["和平区XX街", "沈河区YY路", "铁西区ZZ道"][i], bg_color=bg, text_color=LIGHT_GRAY, font_size=8)
        add_screen_element(slide, screen_left + Inches(3.2), list_y + Inches(0.95)*i + Inches(0.05), Inches(1.1), Inches(0.3),
                           ["远2近1内0瑕0", "远1近0内1瑕0", "远0近0内0瑕0"][i], bg_color=bg, text_color=RGBColor(0x88,0xCC,0x88), font_size=7)
        add_screen_element(slide, screen_left + Inches(0.3), list_y + Inches(0.95)*i + Inches(0.45), Inches(1.5), Inches(0.3),
                           "[拍照]", bg_color=ACCENT_BLUE, font_size=9)
        add_screen_element(slide, screen_left + Inches(1.9), list_y + Inches(0.95)*i + Inches(0.45), Inches(2.0), Inches(0.3),
                           "[查看已拍]", bg_color=RGBColor(0x55,0x55,0x66), font_size=9)
    
    add_page_number(slide, 6, total_slides)
    
    # ==================== 第7页：第五章 拍照功能 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_bar(slide, "第五章  拍照功能")
    
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(6.5), Inches(0.5),
                "拍照流程说明", font_size=24, bold=True, color=ACCENT_BLUE)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(6.5), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "1. 选择拍照类型"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = '微软雅黑'
    add_paragraph(tf, "   点击拍照按钮后，弹出类型选择窗口：远景、近景、内部、瑕疵", font_size=17, color=LIGHT_GRAY, space_before=Pt(4))
    
    add_paragraph(tf, "2. 连续拍摄模式（手动快门）", font_size=18, bold=True, color=WHITE, space_before=Pt(14))
    add_paragraph(tf, "   选择类型后启动系统相机，每按一次快门拍摄一张。拍完一张后相机自动重新调起，", font_size=17, color=LIGHT_GRAY, space_before=Pt(4))
    add_paragraph(tf, "   无需重新选择类型即可拍摄下一张同类型照片。按手机返回键结束拍摄。", font_size=17, color=LIGHT_GRAY, space_before=Pt(2))
    
    add_paragraph(tf, "3. 自动编号", font_size=18, bold=True, color=WHITE, space_before=Pt(14))
    add_paragraph(tf, "   同一类型照片自动编号：01、02、03...", font_size=17, color=LIGHT_GRAY, space_before=Pt(4))
    
    add_paragraph(tf, "4. 水印与命名", font_size=18, bold=True, color=WHITE, space_before=Pt(14))
    add_paragraph(tf, "   拍摄时自动添加水印，并按设置的规则自动命名文件", font_size=17, color=LIGHT_GRAY, space_before=Pt(4))
    
    add_paragraph(tf, "5. 结束拍摄", font_size=18, bold=True, color=WHITE, space_before=Pt(14))
    add_paragraph(tf, "   按手机返回键结束当前类型拍摄，返回客户列表", font_size=17, color=LIGHT_GRAY, space_before=Pt(4))
    
    screen_left = Inches(7.5)
    screen_top = Inches(1.5)
    screen_w = Inches(4.8)
    screen_h = Inches(5.4)
    phone = add_phone_screen(slide, screen_left, screen_top, screen_w, screen_h, "拍照界面")
    
    add_screen_element(slide, screen_left + Inches(0.2), screen_top + Inches(0.5), screen_w - Inches(0.4), Inches(0.4),
                       " <- 返回   张三 - 拍照", bg_color=ACCENT_BLUE, font_size=10)
    
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(1.2), screen_w - Inches(0.6), Inches(0.45),
                       "请选择拍照类型：", font_size=10, bg_color=SCREEN_BG)
    
    types = ["远景", "近景", "内部", "瑕疵"]
    for i, t in enumerate(types):
        ty_x = screen_left + Inches(0.3) + (screen_w - Inches(0.9)) / 4 * i + Inches(0.05)
        ty_w = (screen_w - Inches(0.9)) / 4 - Inches(0.1)
        add_screen_element(slide, ty_x, screen_top + Inches(1.8), ty_w, Inches(0.6),
                           t, bg_color=ACCENT_BLUE if i == 0 else DARK_GRAY, font_size=12)
    
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(2.8), screen_w - Inches(0.6), Inches(2.0),
                       " [ 相机取景框区域 ] \n\n 当前：远景 02\n\n [快门按钮]", font_size=10)
    
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(5.0), screen_w - Inches(0.6), Inches(0.3),
                       "水印预览：2026-06-28 14:30 | 沈阳市和平区...", font_size=7, bg_color=RGBColor(0x33,0x33,0x40))
    
    add_page_number(slide, 7, total_slides)
    
    # ==================== 第8页：第六章 照片查看与管理 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_bar(slide, "第六章  照片查看与管理")
    
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(6.5), Inches(0.5),
                "照片管理功能", font_size=24, bold=True, color=ACCENT_BLUE)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(6.5), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "查看已拍照片"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = '微软雅黑'
    add_paragraph(tf, "点击「查看已拍」按钮，进入照片浏览页面，按类型（远景/近景/内部/瑕疵）分类展示该客户所有已拍照片，以缩略图网格形式排列。", font_size=17, color=LIGHT_GRAY, space_before=Pt(6))
    
    add_paragraph(tf, "删除单张照片", font_size=20, bold=True, color=WHITE, space_before=Pt(16))
    add_paragraph(tf, "点击单张照片下方的「删除」按钮，弹出确认对话框，需二次确认后方可删除，防止误操作。", font_size=17, color=LIGHT_GRAY, space_before=Pt(6))
    
    add_paragraph(tf, "删除全部照片", font_size=20, bold=True, color=WHITE, space_before=Pt(16))
    add_paragraph(tf, "页面提供「删除全部」按钮，点击后同样需要二次确认，可快速清空该客户所有照片后重新拍摄。", font_size=17, color=LIGHT_GRAY, space_before=Pt(6))
    
    screen_left = Inches(7.5)
    screen_top = Inches(1.5)
    screen_w = Inches(4.8)
    screen_h = Inches(5.4)
    phone = add_phone_screen(slide, screen_left, screen_top, screen_w, screen_h, "照片查看页")
    
    add_screen_element(slide, screen_left + Inches(0.2), screen_top + Inches(0.5), screen_w - Inches(0.4), Inches(0.4),
                       " <- 返回   张三 - 已拍照片", bg_color=ACCENT_BLUE, font_size=9)
    
    add_screen_element(slide, screen_left + Inches(0.2), screen_top + Inches(1.0), screen_w - Inches(0.4), Inches(0.35),
                       " [远景(3)] [近景(2)] [内部(0)] [瑕疵(1)]", bg_color=DARK_GRAY, font_size=8)
    
    thumb_y = screen_top + Inches(1.5)
    thumb_w = (screen_w - Inches(0.6)) / 3
    for r in range(3):
        for c in range(3):
            idx = r * 3 + c
            if idx < 6:
                add_screen_element(slide, screen_left + Inches(0.2) + thumb_w * c + Inches(0.03),
                                   thumb_y + Inches(0.9) * r + Inches(0.03),
                                   thumb_w - Inches(0.06), Inches(0.85),
                                   f"照片{idx+1:02d}", bg_color=RGBColor(0x40,0x40,0x50), font_size=9)
    
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(4.5), screen_w - Inches(0.6), Inches(0.45),
                       "[删除全部]（需二次确认）", bg_color=RGBColor(0x8B,0x3A,0x3A), font_size=10)
    
    add_page_number(slide, 8, total_slides)
    
    # ==================== 第9页：第七章 设置页面 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_bar(slide, "第七章  设置页面")
    
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(6.5), Inches(0.5),
                "设置项说明", font_size=24, bold=True, color=ACCENT_BLUE)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(6.5), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    settings = [
        ("返回按钮：", "点击返回主界面，设置不保存则提示"),
        ("命名规则：", "4段下拉选择器，组合照片文件命名方式（详见第九章）"),
        ("水印设置：", "3段内容选择 + 字号大小选择 + 水印位置选择"),
        ("  - 水印段1/2/3：", "选择水印显示的内容字段"),
        ("  - 字号：", "大 / 中 / 小 三档可选"),
        ("  - 位置：", "右上 / 右下 / 左上 / 左下 四个角位置"),
        ("保存按钮：", "点击后保存所有设置，立即生效"),
    ]
    
    p = tf.paragraphs[0]
    p.text = settings[0][0] + settings[0][1]
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = '微软雅黑'
    for title, desc in settings[1:]:
        p = add_paragraph(tf, title + desc, font_size=17, color=LIGHT_GRAY, space_before=Pt(10))
        if title.startswith("  - "):
            p.runs[0].font.bold = False
        else:
            p.runs[0].font.bold = True
    
    screen_left = Inches(7.5)
    screen_top = Inches(1.5)
    screen_w = Inches(4.8)
    screen_h = Inches(5.4)
    phone = add_phone_screen(slide, screen_left, screen_top, screen_w, screen_h, "设置页面")
    
    add_screen_element(slide, screen_left + Inches(0.2), screen_top + Inches(0.5), screen_w - Inches(0.4), Inches(0.4),
                       " <- 返回   设置", bg_color=ACCENT_BLUE, font_size=11)
    
    set_y = screen_top + Inches(1.1)
    add_screen_element(slide, screen_left + Inches(0.3), set_y, Inches(1.5), Inches(0.35),
                       "命名规则：", font_size=9, bg_color=SCREEN_BG)
    for i in range(4):
        add_screen_element(slide, screen_left + Inches(1.9) + Inches(0.65)*i, set_y, Inches(0.6), Inches(0.35),
                           f"段{i+1}v", bg_color=DARK_GRAY, font_size=7)
    
    add_screen_element(slide, screen_left + Inches(0.3), set_y + Inches(0.6), Inches(1.5), Inches(0.35),
                       "水印内容：", font_size=9, bg_color=SCREEN_BG)
    for i in range(3):
        add_screen_element(slide, screen_left + Inches(1.9) + Inches(0.85)*i, set_y + Inches(0.6), Inches(0.8), Inches(0.35),
                           f"水{i+1}v", bg_color=DARK_GRAY, font_size=7)
    
    add_screen_element(slide, screen_left + Inches(0.3), set_y + Inches(1.2), Inches(1.5), Inches(0.35),
                       "水印字号：", font_size=9, bg_color=SCREEN_BG)
    add_screen_element(slide, screen_left + Inches(1.9), set_y + Inches(1.2), Inches(2.2), Inches(0.35),
                       "大  |  中  |  小", font_size=9)
    
    add_screen_element(slide, screen_left + Inches(0.3), set_y + Inches(1.8), Inches(1.5), Inches(0.35),
                       "水印位置：", font_size=9, bg_color=SCREEN_BG)
    add_screen_element(slide, screen_left + Inches(1.9), set_y + Inches(1.8), Inches(2.2), Inches(0.35),
                       "左上 右上 左下 右下", font_size=8)
    
    add_screen_element(slide, screen_left + Inches(0.8), screen_top + Inches(4.5), screen_w - Inches(1.6), Inches(0.6),
                       "[ 保 存 设 置 ]", bg_color=ACCENT_BLUE, font_size=14)
    
    add_page_number(slide, 9, total_slides)
    
    # ==================== 第10页：第八章 水印说明 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_bar(slide, "第八章  水印说明")
    
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(6.5), Inches(0.5),
                "水印配置详解", font_size=24, bold=True, color=ACCENT_BLUE)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(6.5), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    add_paragraph_ex = lambda txt, sz=17, b=False, c=LIGHT_GRAY, sb=Pt(8): add_paragraph(tf, txt, font_size=sz, bold=b, color=c, space_before=sb)
    
    p = tf.paragraphs[0]
    p.text = "水印内容段选择"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = '微软雅黑'
    add_paragraph_ex("每段水印内容可从以下选项中选择：", sb=Pt(8))
    for opt in ["拍摄时间（精确到分钟）", "客户名称", "地址全（概+精确拼接）", "地址概", "地址精确", "GPS定位坐标", "空值（不显示该段）"]:
        add_paragraph_ex("  - " + opt, sb=Pt(4))
    
    add_paragraph_ex("字号选择", sz=20, b=True, c=WHITE, sb=Pt(16))
    add_paragraph_ex("  - 大字号：约48pt，适合远距离查看", sb=Pt(4))
    add_paragraph_ex("  - 中字号：约36pt，默认推荐大小", sb=Pt(4))
    add_paragraph_ex("  - 小字号：约24pt，适合信息较多时使用", sb=Pt(4))
    
    add_paragraph_ex("位置选择", sz=20, b=True, c=WHITE, sb=Pt(16))
    add_paragraph_ex("  - 右上角：不遮挡主体时的首选位置", sb=Pt(4))
    add_paragraph_ex("  - 右下角：常用位置，适合横向拍摄", sb=Pt(4))
    add_paragraph_ex("  - 左上角：天空/浅色背景时使用", sb=Pt(4))
    add_paragraph_ex("  - 左下角：地面/深色背景时使用", sb=Pt(4))
    
    screen_left = Inches(7.5)
    screen_top = Inches(1.5)
    screen_w = Inches(4.8)
    screen_h = Inches(5.4)
    phone = add_phone_screen(slide, screen_left, screen_top, screen_w, screen_h, "水印位置示意")
    
    photo_area = add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(0.5), screen_w - Inches(0.6), Inches(4.6),
                                    "[ 示例照片区域 ]", bg_color=RGBColor(0x35,0x35,0x42), font_size=10)
    
    add_screen_element(slide, screen_left + Inches(0.5), screen_top + Inches(0.7), Inches(1.6), Inches(0.4),
                       "【左上水印】", bg_color=RGBColor(0x55,0x55,0x55), font_size=7)
    add_screen_element(slide, screen_left + screen_w - Inches(2.1), screen_top + Inches(0.7), Inches(1.6), Inches(0.4),
                       "【右上水印】", bg_color=RGBColor(0x55,0x55,0x55), font_size=7)
    add_screen_element(slide, screen_left + Inches(0.5), screen_top + Inches(4.4), Inches(1.6), Inches(0.4),
                       "【左下水印】", bg_color=RGBColor(0x55,0x55,0x55), font_size=7)
    add_screen_element(slide, screen_left + screen_w - Inches(2.1), screen_top + Inches(4.4), Inches(1.6), Inches(0.4),
                       "【右下水印】\n2026-06-28\n张三 和平区...", bg_color=RGBColor(0x55,0x55,0x55), font_size=7)
    
    add_page_number(slide, 10, total_slides)
    
    # ==================== 第11页：第九章 文件命名规则 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_bar(slide, "第九章  文件命名规则")
    
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(6.5), Inches(0.5),
                "命名规则说明", font_size=24, bold=True, color=ACCENT_BLUE)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(6.5), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "4段组合命名"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = '微软雅黑'
    add_paragraph(tf, "文件名由4个可配置的字段段组合而成，各段之间用下划线「_」连接。每段可选择以下字段之一：", font_size=17, color=LIGHT_GRAY, space_before=Pt(8))
    
    fields = [
        "拍摄时间（格式：20260628_1430）",
        "客户名",
        "地址全（地址概+地址精确）",
        "地址概",
        "地址精确",
        "空值（该段省略）",
    ]
    for f in fields:
        add_paragraph(tf, "  - " + f, font_size=17, color=LIGHT_GRAY, space_before=Pt(5))
    
    add_paragraph(tf, "自动追加后缀", font_size=20, bold=True, color=WHITE, space_before=Pt(16))
    add_paragraph(tf, "在4段组合的文件名之后，系统会自动追加：拍照类型（远景/近景/内部/瑕疵）和照片编号（01/02/...）。", font_size=17, color=LIGHT_GRAY, space_before=Pt(8))
    
    add_paragraph(tf, "命名示例", font_size=20, bold=True, color=WHITE, space_before=Pt(16))
    add_paragraph(tf, "  20260628_1430_张三_沈阳市和平区XX街123号_-远景-01.jpg", font_size=15, color=ACCENT_BLUE, space_before=Pt(6))
    
    screen_left = Inches(7.5)
    screen_top = Inches(1.5)
    screen_w = Inches(4.8)
    screen_h = Inches(5.4)
    phone = add_phone_screen(slide, screen_left, screen_top, screen_w, screen_h, "命名规则示例")
    
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(0.6), screen_w - Inches(0.6), Inches(0.4),
                       "文件名结构：", bg_color=SCREEN_BG, font_size=11)
    
    parts_y = screen_top + Inches(1.2)
    parts = ["段1\n拍摄时间", "段2\n客户名", "段3\n地址概", "段4\n地址精确"]
    part_w = (screen_w - Inches(0.8)) / 4
    for i, pt in enumerate(parts):
        c = ACCENT_BLUE if i % 2 == 0 else RGBColor(0x3A,0x6A,0xA0)
        add_screen_element(slide, screen_left + Inches(0.4) + part_w * i, parts_y, part_w - Inches(0.05), Inches(0.8),
                           pt, bg_color=c, font_size=8)
    
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(2.3), Inches(0.15), Inches(0.15), "+", font_size=12, bg_color=SCREEN_BG, text_color=LIGHT_GRAY)
    add_screen_element(slide, screen_left + Inches(0.3) + part_w, screen_top + Inches(2.3), Inches(0.15), Inches(0.15), "+", font_size=12, bg_color=SCREEN_BG, text_color=LIGHT_GRAY)
    add_screen_element(slide, screen_left + Inches(0.3) + part_w*2, screen_top + Inches(2.3), Inches(0.15), Inches(0.15), "+", font_size=12, bg_color=SCREEN_BG, text_color=LIGHT_GRAY)
    add_screen_element(slide, screen_left + Inches(0.3) + part_w*3, screen_top + Inches(2.3), Inches(0.15), Inches(0.15), "+", font_size=12, bg_color=SCREEN_BG, text_color=LIGHT_GRAY)
    
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(2.7), screen_w - Inches(0.6), Inches(0.5),
                       "+ 类型（远景/近景/内部/瑕疵）", bg_color=RGBColor(0x66,0x55,0x33), font_size=10)
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(3.4), screen_w - Inches(0.6), Inches(0.5),
                       "+ 编号（01/02/03...）", bg_color=RGBColor(0x33,0x55,0x33), font_size=10)
    
    add_screen_element(slide, screen_left + Inches(0.3), screen_top + Inches(4.2), screen_w - Inches(0.6), Inches(0.8),
                       "最终文件名：\n20260628_1430_张三_和平区_-远景-01.jpg", bg_color=RGBColor(0x25,0x45,0x25), font_size=8)
    
    add_page_number(slide, 11, total_slides)
    
    # ==================== 第12页：第十章 常见问题 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_bar(slide, "第十章  常见问题")
    
    faqs = [
        ("Q1: 点击拍照后相机没有启动？", 
         "A: 请先确认已授予相机权限（手机设置 -> 应用管理 -> 本应用 -> 权限管理）。应用采用三重兼容策略调用系统相机，如仍无法启动，请确认手机系统相机应用是否正常工作。部分定制系统可能需要手动在权限管理中开启「相机」权限。"),
        ("Q2: GPS定位不准或无法获取？", 
         "A: 请确保手机已开启GPS定位功能，并在开阔地带使用。室内可能影响GPS信号接收。水印中定位信息需要等待GPS定位成功后才会显示，首次定位可能需要30秒左右。"),
        ("Q3: 照片保存在哪里？", 
         "A: 照片保存在应用私有目录的 LoanPhotos/ 文件夹下，按客户名分子文件夹存放。系统相册也可查看。导出时可通过手机文件管理器或USB连接电脑拷贝。"),
        ("Q4: 如何结束连续拍摄？", 
         "A: 在拍摄状态下按手机返回键即可结束当前类型的拍摄，返回客户列表。每次按快门拍摄一张后，相机会自动重新调起等待下一张，无需重新选择类型。"),
        ("Q5: 删除照片能否恢复？", 
         "A: 不能。删除照片前会弹出二次确认对话框，确认后照片将被永久删除且无法恢复，请谨慎操作。"),
    ]
    
    y_pos = Inches(1.5)
    for q, a in faqs:
        qbox = add_textbox(slide, Inches(0.5), y_pos, Inches(12.3), Inches(0.4),
                           q, font_size=18, bold=True, color=ACCENT_BLUE)
        atxBox = slide.shapes.add_textbox(Inches(0.8), y_pos + Inches(0.4), Inches(11.8), Inches(0.7))
        atf = atxBox.text_frame
        atf.word_wrap = True
        ap = atf.paragraphs[0]
        ap.text = a
        ap.font.size = Pt(15)
        ap.font.color.rgb = LIGHT_GRAY
        ap.font.name = '微软雅黑'
        y_pos += Inches(1.1)
    
    add_page_number(slide, 12, total_slides)
    
    # ==================== 第13页：封底 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    for i in range(5):
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(8 + i * 0.15), Inches(2.5 + i * 0.08),
            Inches(4 - i * 0.3), Inches(0.04)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = ACCENT_BLUE if i == 0 else RGBColor(0x3A + i*5, 0x70 + i*5, 0xB0 + i*5)
        accent_line.line.fill.background()
    
    add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(0.8),
                "感谢使用", font_size=40, bold=True, color=WHITE)
    
    add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
                "如有问题或建议，请联系开发者", font_size=22, color=LIGHT_GRAY)
    
    contact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.3), Inches(6), Inches(1.8))
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = DARK_GRAY
    contact_box.line.color.rgb = ACCENT_BLUE
    contact_box.line.width = Pt(1)
    
    ctf = contact_box.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Inches(0.3)
    ctf.margin_top = Inches(0.2)
    cp = ctf.paragraphs[0]
    cp.text = "联系方式"
    cp.font.size = Pt(20)
    cp.font.bold = True
    cp.font.color.rgb = ACCENT_BLUE
    cp.font.name = '微软雅黑'
    
    add_paragraph(ctf, "作者：王硕", font_size=18, color=WHITE, space_before=Pt(12))
    add_paragraph(ctf, "电话：15940454123", font_size=18, color=WHITE, space_before=Pt(8))
    
    add_textbox(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                "资产盘点专项拍照工具 v3.6.0", font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    
    add_page_number(slide, 13, total_slides)
    
    output_path = r"C:\Users\Administrator\Desktop\资产盘点拍照工具-使用说明书.pptx"
    prs.save(output_path)
    print(f"PPT已成功生成：{output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
